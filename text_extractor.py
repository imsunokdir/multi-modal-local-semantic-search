import os
import io
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pandas as pd
import pytesseract
from pdf2image import convert_from_path
from PIL import Image

# ✅ CORRECT - pointing to the .exe file
pytesseract.pytesseract.tesseract_cmd = r"C:\Users\tangi\AppData\Local\Programs\Tesseract-OCR\tesseract.exe"

def extract_text_from_file(filepath):  # ✅ Renamed function
    ext = filepath.lower().split('.')[-1]

    # -------- TXT --------
    if ext == "txt":
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        except:
            return ""

    # -------- PDF --------
    elif ext == "pdf":
        text = ""
        try:
            doc = fitz.open(filepath)
            for page in doc:
                page_text = page.get_text()
                text += page_text

            # If no text found, use OCR
            if len(text.strip()) == 0:
                print(f"🖼️ Using OCR for: {os.path.basename(filepath)}")
                for page_num in range(len(doc)):
                    pix = doc[page_num].get_pixmap(dpi=200, alpha=False)
                    img_bytes = pix.tobytes("png")
                    img = Image.open(io.BytesIO(img_bytes))
                    text += pytesseract.image_to_string(img)
            
            doc.close()
            return text
            
        except PermissionError as e:
            print(f"⚠️ Permission denied for PDF {filepath}: {e}")
            return ""
        except Exception as e:
            print(f"⚠️ Error reading PDF {filepath}: {e}")
            return ""

    # -------- DOCX --------
    elif ext == "docx":
        try:
            doc = docx.Document(filepath)
            return "\n".join([p.text for p in doc.paragraphs])
        except:
            return ""

    # -------- PPTX --------
    elif ext == "pptx":
        try:
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except:
            return ""

    # -------- CSV --------
    elif ext == "csv":
        try:
            df = pd.read_csv(filepath)
            return df.to_string()
        except:
            return ""

    # -------- UNKNOWN FORMAT --------
    return ""